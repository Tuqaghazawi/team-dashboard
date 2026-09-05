"""Generating MDC and planning-round slide decks.

The layout follows the decks the teams already present: one slide per patient,
titled with the patient's name and MRN, then the history line, the case line,
each investigation with its report, any treatment already given, and the
decision line at the bottom.

Two things sit at the top of every slide because the room asks for them first:
the MRN and the genetic testing status.

Guideline evidence goes into the slide's notes field, never onto the slide
itself — it is there to support the discussion, not to pre-empt it.
"""

from io import BytesIO

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

from patients.models import Investigation

# House colours, matching the dashboard.
TEAL = RGBColor(0x02, 0x80, 0x90)
INK = RGBColor(0x13, 0x34, 0x3B)
MUTED = RGBColor(0x5B, 0x7B, 0x7A)
GOLD = RGBColor(0xE9, 0xA2, 0x3B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LINE = RGBColor(0xE2, 0xEB, 0xEA)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# The order investigations read on a slide — endoscopy, tissue, then imaging.
KIND_ORDER = [
    Investigation.Kind.COLONOSCOPY,
    Investigation.Kind.SIGMOIDOSCOPY,
    Investigation.Kind.GASTROSCOPY,
    Investigation.Kind.MAMMOGRAM,
    Investigation.Kind.BREAST_US,
    Investigation.Kind.NECK_US,
    Investigation.Kind.FNA,
    Investigation.Kind.PATHOLOGY,
    Investigation.Kind.CEA,
    Investigation.Kind.TUMOR_MARKERS,
    Investigation.Kind.THYROID_FUNCTION,
    Investigation.Kind.CAP_CT,
    Investigation.Kind.PELVIC_MRI,
    Investigation.Kind.ABDOMEN_MRI,
    Investigation.Kind.BREAST_MRI,
    Investigation.Kind.LOCAL_MRI,
    Investigation.Kind.PET_CT,
    Investigation.Kind.BONE_SCAN,
    Investigation.Kind.ECHO,
    Investigation.Kind.PFT,
    Investigation.Kind.GENETICS,
]


def build_mdc_deck(mdc_name, meeting_date, listings, presenter="", evidence=None):
    """A deck for one MDC meeting. ``listings`` are MDCListing objects."""
    prs = _new_presentation()
    _title_slide(prs, f"{mdc_name} MDC", presenter, meeting_date)
    for listing in listings:
        _case_slide(
            prs,
            listing.patient,
            decision_label="MDC",
            decision_text=listing.decision,
            footer=f"{mdc_name} MDC · {meeting_date:%d/%m/%Y}",
            notes=(evidence or {}).get(listing.patient_id, ""),
        )
    return _save(prs)


def build_planning_deck(team, meeting_date, patients, evidence=None):
    """A planning-round deck for one team's operative patients."""
    prs = _new_presentation()
    _title_slide(prs, "Planning", team.consultant, meeting_date)
    for patient in patients:
        booking = next(
            (b for b in patient.surgery_bookings.all() if not b.performed), None
        )
        _case_slide(
            prs,
            patient,
            decision_label="For",
            decision_text=_planned_procedure(patient, booking),
            footer=f"Planning · {team.consultant} · {meeting_date:%d/%m/%Y}",
            notes=(evidence or {}).get(patient.pk, ""),
        )
    return _save(prs)


# --- slide construction -------------------------------------------------------

def _new_presentation():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def _blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _title_slide(prs, title, presenter, meeting_date):
    slide = _blank(prs)
    _rect(slide, 0, 0, SLIDE_W, SLIDE_H, TEAL)

    box = _textbox(slide, Inches(1.0), Inches(2.5), Inches(11.3), Inches(2.5))
    frame = box.text_frame
    frame.word_wrap = True

    para = frame.paragraphs[0]
    _run(para, title, size=48, bold=True, color=WHITE)
    para.alignment = PP_ALIGN.CENTER

    if presenter:
        p2 = frame.add_paragraph()
        _run(p2, presenter, size=24, color=RGBColor(0xD7, 0xEE, 0xF0))
        p2.alignment = PP_ALIGN.CENTER
        p2.space_before = Pt(18)

    p3 = frame.add_paragraph()
    _run(p3, f"{meeting_date:%d/%m/%Y}", size=18, color=RGBColor(0xB8, 0xDE, 0xE2))
    p3.alignment = PP_ALIGN.CENTER
    p3.space_before = Pt(10)


def _case_slide(prs, patient, decision_label, decision_text, footer, notes=""):
    slide = _blank(prs)

    # --- header band: name + MRN, with genetics on the right ---
    _rect(slide, 0, 0, SLIDE_W, Inches(1.0), TEAL)

    head = _textbox(slide, Inches(0.45), Inches(0.12), Inches(8.6), Inches(0.78))
    head.text_frame.word_wrap = True
    hp = head.text_frame.paragraphs[0]
    _run(hp, patient.name, size=26, bold=True, color=WHITE)
    _run(hp, f"   {patient.mrn}", size=20, color=RGBColor(0xCD, 0xEA, 0xEA))

    genetics = patient.genetic_testing.strip() or "Genetics: not recorded"
    if not genetics.lower().startswith("genetic"):
        genetics = f"Genetics: {genetics}"
    gbox = _textbox(slide, Inches(9.2), Inches(0.24), Inches(3.7), Inches(0.55))
    gbox.text_frame.word_wrap = True
    gp = gbox.text_frame.paragraphs[0]
    gp.alignment = PP_ALIGN.RIGHT
    _run(gp, genetics, size=14, bold=True, color=GOLD)

    # --- body ---
    body = _textbox(slide, Inches(0.5), Inches(1.2), Inches(12.4), Inches(5.6))
    frame = body.text_frame
    frame.word_wrap = True
    frame.vertical_anchor = MSO_ANCHOR.TOP
    first = True

    for label, text, emphasis in _case_lines(patient):
        para = frame.paragraphs[0] if first else frame.add_paragraph()
        first = False
        para.space_after = Pt(5)
        if label:
            _run(para, f"{label}: ", size=14, bold=True, color=TEAL)
        _run(para, text, size=14, bold=emphasis, color=INK)

    # --- decision line ---
    dec = _textbox(slide, Inches(0.5), Inches(6.55), Inches(12.4), Inches(0.6))
    dec.text_frame.word_wrap = True
    dp = dec.text_frame.paragraphs[0]
    _run(dp, f"{decision_label}: ", size=17, bold=True, color=TEAL)
    _run(dp, decision_text or "", size=17, bold=True, color=INK)

    # --- footer ---
    _rect(slide, 0, Inches(7.18), SLIDE_W, Inches(0.32), RGBColor(0xF4, 0xF8, 0xF7))
    foot = _textbox(slide, Inches(0.45), Inches(7.16), Inches(12.4), Inches(0.3))
    fp = foot.text_frame.paragraphs[0]
    _run(fp, footer, size=10, color=MUTED)

    if notes:
        slide.notes_slide.notes_text_frame.text = notes
    return slide


def _case_lines(patient):
    """The body of a case slide as (label, text, emphasis) rows."""
    lines = []

    # History line — "73-year-old female patient, medically free, PS: 0."
    bits = [f"{patient.age}-year-old"]
    if patient.sex:
        bits.append(patient.get_sex_display().lower())
    bits.append("patient")
    history = " ".join(bits)
    if patient.comorbidities:
        history += f", {patient.comorbidities}"
    lines.append(("", history, False))

    if patient.family_history:
        lines.append(("Family history", patient.family_history, False))

    case = f"Case of {patient.diagnosis}"
    if patient.clinical_stage:
        case += f" [{patient.clinical_stage}]"
    lines.append(("", case, True))

    baseline = _ordered(patient, Investigation.Purpose.BASELINE)
    for item in baseline:
        if item.kind == Investigation.Kind.GENETICS:
            continue  # already in the header
        lines.append((item.get_kind_display(), _result_or_status(item), False))

    # Treatment already given, then restaging.
    for course in patient.treatment_courses.all():
        lines.append((
            "S/P",
            f"{course.get_kind_display()} — {course.regimen} "
            f"({course.completed_cycles}/{course.total_cycles} cycles)",
            True,
        ))

    restaging = _ordered(patient, Investigation.Purpose.RESTAGING)
    if restaging:
        lines.append(("", "Restaging:", True))
        for item in restaging:
            lines.append((item.get_kind_display(), _result_or_status(item), False))

    # Operative detail, for post-op re-discussion.
    for booking in patient.surgery_bookings.all():
        if booking.performed:
            lines.append(("S/P", f"{booking.procedure} — {booking.performed_on}", True))
            if booking.final_pathology:
                lines.append(("Final pathology", booking.final_pathology, False))

    return lines


def _ordered(patient, purpose):
    """This patient's investigations for a purpose, in presentation order."""
    items = [i for i in patient.investigations.all() if i.purpose == purpose]
    position = {kind: n for n, kind in enumerate(KIND_ORDER)}
    return sorted(items, key=lambda i: position.get(i.kind, len(KIND_ORDER)))


def _result_or_status(item):
    if item.status == Investigation.Status.READY and item.result_text:
        return item.result_text
    return f"[{item.get_status_display().lower()}]"


def _planned_procedure(patient, booking):
    if booking:
        return booking.procedure
    listing = patient.mdc_listings.filter(decision_category="SURGERY").first()
    return listing.decision if listing else ""


# --- small pptx helpers -------------------------------------------------------

def _rect(slide, left, top, width, height, colour):
    from pptx.enum.shapes import MSO_SHAPE

    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(int(left)), Emu(int(top)),
                                   Emu(int(width)), Emu(int(height)))
    shape.fill.solid()
    shape.fill.fore_color.rgb = colour
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def _textbox(slide, left, top, width, height):
    box = slide.shapes.add_textbox(left, top, width, height)
    box.text_frame.margin_left = 0
    box.text_frame.margin_right = 0
    box.text_frame.margin_top = 0
    box.text_frame.margin_bottom = 0
    return box


def _run(paragraph, text, size=14, bold=False, color=INK):
    run = paragraph.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Segoe UI"
    return run


def _save(prs):
    stream = BytesIO()
    prs.save(stream)
    stream.seek(0)
    return stream
